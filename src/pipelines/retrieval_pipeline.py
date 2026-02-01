"""
Agentic RAG System - Main Implementation
=========================================

Two Pipeline Architecture:
1. Data Pipeline: Extraction → Chunking → Metadata → Embedding → Indexing
2. Retrieval Pipeline: Query → Hybrid Search → Agentic Re-retrieval → LLM Answer

Tech Stack:
- Embeddings: all-mpnet-base-v2 (sentence-transformers)
- LLM: Ollama (llama3.2 - self-hosted, free)
- Vector DB: FAISS (local, fast)
"""

import os
import json
import numpy as np
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
import re
import math
from collections import Counter

# Document processing imports
import PyPDF2
import docx
import openpyxl
from pptx import Presentation

# ML imports
from sentence_transformers import SentenceTransformer
import faiss

# LLM import
import requests


# ============================================================================
# DATA STRUCTURES
# ============================================================================

@dataclass
class DocumentMetadata:
    """Metadata for each document chunk"""
    chunk_id: str
    filename: str
    file_type: str
    page: Optional[int] = None
    sheet: Optional[str] = None
    slide: Optional[int] = None
    section: Optional[str] = None
    timestamp: str = field(default_factory=lambda: datetime.now().isoformat())
    chunk_index: int = 0
    total_chunks: int = 0


@dataclass
class DocumentChunk:
    """A chunk of document with metadata and embedding"""
    content: str
    metadata: DocumentMetadata
    embedding: Optional[np.ndarray] = None


@dataclass
class RetrievalResult:
    """Result from hybrid retrieval"""
    chunk: DocumentChunk
    bm25_score: float
    vector_score: float
    combined_score: float


# ============================================================================
# PIPELINE 1: DATA PROCESSING PIPELINE
# ============================================================================

class DocumentExtractor:
    """Extract text from various document formats"""
    
    @staticmethod
    def extract_pdf(filepath: str) -> List[Tuple[str, int]]:
        """Extract text from PDF - returns (text, page_number) tuples"""
        pages = []
        with open(filepath, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for i, page in enumerate(pdf_reader.pages):
                text = page.extract_text()
                if text.strip():
                    pages.append((text, i + 1))
        return pages
    
    @staticmethod
    def extract_docx(filepath: str) -> str:
        """Extract text from DOCX"""
        doc = docx.Document(filepath)
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        return '\n'.join(paragraphs)
    
    @staticmethod
    def extract_pptx(filepath: str) -> List[Tuple[str, int]]:
        """Extract text from PPTX - returns (text, slide_number) tuples"""
        prs = Presentation(filepath)
        slides = []
        for i, slide in enumerate(prs.slides):
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
            if text_parts:
                slides.append(('\n'.join(text_parts), i + 1))
        return slides
    
    @staticmethod
    def extract_excel(filepath: str) -> List[Tuple[str, str]]:
        """Extract text from Excel - returns (text, sheet_name) tuples"""
        wb = openpyxl.load_workbook(filepath, data_only=True)
        sheets = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_text = ' | '.join([str(cell) if cell is not None else '' for cell in row])
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                sheets.append(('\n'.join(rows), sheet_name))
        return sheets
    
    @staticmethod
    def extract_txt(filepath: str) -> str:
        """Extract text from TXT file"""
        with open(filepath, 'r', encoding='utf-8') as file:
            return file.read()


class TextChunker:
    """Chunk text with overlap for better context preservation"""
    
    def __init__(self, chunk_size: int = 512, overlap: int = 50):
        self.chunk_size = chunk_size
        self.overlap = overlap
    
    def chunk_text(self, text: str) -> List[str]:
        """Split text into overlapping chunks"""
        chunks = []
        text = self._clean_text(text)
        
        # Split by sentences for better boundaries
        sentences = re.split(r'(?<=[.!?])\s+', text)
        
        current_chunk = []
        current_size = 0
        
        for sentence in sentences:
            sentence_size = len(sentence)
            
            if current_size + sentence_size > self.chunk_size and current_chunk:
                # Save current chunk
                chunks.append(' '.join(current_chunk))
                
                # Start new chunk with overlap
                overlap_text = ' '.join(current_chunk[-2:]) if len(current_chunk) >= 2 else ''
                current_chunk = [overlap_text] if overlap_text else []
                current_size = len(overlap_text)
            
            current_chunk.append(sentence)
            current_size += sentence_size
        
        # Add final chunk
        if current_chunk:
            chunks.append(' '.join(current_chunk))
        
        return chunks
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize text"""
        # Remove extra whitespace
        text = re.sub(r'\s+', ' ', text)
        # Remove control characters
        text = ''.join(char for char in text if ord(char) >= 32 or char in '\n\t')
        return text.strip()


class BM25Index:
    """BM25 index for sparse retrieval"""
    
    def __init__(self, k1: float = 1.5, b: float = 0.75):
        self.k1 = k1
        self.b = b
        self.chunks: List[DocumentChunk] = []
        self.doc_lengths: List[int] = []
        self.avg_doc_length: float = 0.0
        self.idf_scores: Dict[str, float] = {}
    
    def index(self, chunks: List[DocumentChunk]):
        """Build BM25 index"""
        self.chunks = chunks
        self.doc_lengths = []
        term_doc_freq = {}
        
        # Calculate document frequencies
        for chunk in chunks:
            tokens = self._tokenize(chunk.content)
            self.doc_lengths.append(len(tokens))
            
            unique_terms = set(tokens)
            for term in unique_terms:
                term_doc_freq[term] = term_doc_freq.get(term, 0) + 1
        
        self.avg_doc_length = sum(self.doc_lengths) / len(self.doc_lengths) if self.doc_lengths else 0
        
        # Calculate IDF scores
        num_docs = len(chunks)
        for term, df in term_doc_freq.items():
            idf = math.log((num_docs - df + 0.5) / (df + 0.5) + 1.0)
            self.idf_scores[term] = idf
    
    def search(self, query: str, top_k: int = 50) -> List[Tuple[DocumentChunk, float]]:
        """Search using BM25"""
        query_tokens = self._tokenize(query)
        scores = []
        
        for i, chunk in enumerate(self.chunks):
            doc_tokens = self._tokenize(chunk.content)
            score = self._calculate_bm25(query_tokens, doc_tokens, self.doc_lengths[i])
            scores.append((chunk, score))
        
        scores.sort(key=lambda x: x[1], reverse=True)
        return scores[:top_k]
    
    def _tokenize(self, text: str) -> List[str]:
        """Simple tokenization"""
        return re.findall(r'\b\w+\b', text.lower())
    
    def _calculate_bm25(self, query_tokens: List[str], doc_tokens: List[str], doc_length: int) -> float:
        """Calculate BM25 score"""
        score = 0.0
        doc_term_freq = Counter(doc_tokens)
        
        for term in query_tokens:
            if term not in doc_term_freq:
                continue
            
            tf = doc_term_freq[term]
            idf = self.idf_scores.get(term, 0)
            
            numerator = tf * (self.k1 + 1)
            denominator = tf + self.k1 * (1 - self.b + self.b * (doc_length / self.avg_doc_length))
            score += idf * (numerator / denominator)
        
        return score


class DataProcessingPipeline:
    """
    PIPELINE 1: Data Processing
    Handles: Extraction → Chunking → Metadata → Embedding → Indexing
    """
    
    def __init__(self, embedding_model_name: str = "sentence-transformers/all-mpnet-base-v2"):
        print(f"Initializing Data Processing Pipeline...")
        print(f"Loading embedding model: {embedding_model_name}")
        
        self.embedding_model = SentenceTransformer(embedding_model_name)
        self.chunker = TextChunker(chunk_size=512, overlap=50)
        self.extractor = DocumentExtractor()
        
        # Indexes
        self.bm25_index = BM25Index()
        self.vector_index = None
        self.chunks: List[DocumentChunk] = []
        
        print("✓ Data Processing Pipeline ready")
    
    def process_document(self, filepath: str) -> List[DocumentChunk]:
        """Process a single document through the pipeline"""
        print(f"\nProcessing: {filepath}")
        
        filename = os.path.basename(filepath)
        file_ext = os.path.splitext(filename)[1].lower()
        
        # STAGE 1: Extract text
        print("  [1/5] Extracting text...")
        chunks = []
        
        if file_ext == '.pdf':
            pages = self.extractor.extract_pdf(filepath)
            for text, page_num in pages:
                text_chunks = self.chunker.chunk_text(text)
                for i, chunk_text in enumerate(text_chunks):
                    metadata = DocumentMetadata(
                        chunk_id=f"{filename}_p{page_num}_c{i}",
                        filename=filename,
                        file_type="PDF",
                        page=page_num,
                        chunk_index=i,
                        total_chunks=len(text_chunks)
                    )
                    chunks.append(DocumentChunk(content=chunk_text, metadata=metadata))
        
        elif file_ext == '.docx':
            text = self.extractor.extract_docx(filepath)
            text_chunks = self.chunker.chunk_text(text)
            for i, chunk_text in enumerate(text_chunks):
                metadata = DocumentMetadata(
                    chunk_id=f"{filename}_c{i}",
                    filename=filename,
                    file_type="DOCX",
                    chunk_index=i,
                    total_chunks=len(text_chunks)
                )
                chunks.append(DocumentChunk(content=chunk_text, metadata=metadata))
        
        elif file_ext == '.pptx':
            slides = self.extractor.extract_pptx(filepath)
            for text, slide_num in slides:
                text_chunks = self.chunker.chunk_text(text)
                for i, chunk_text in enumerate(text_chunks):
                    metadata = DocumentMetadata(
                        chunk_id=f"{filename}_s{slide_num}_c{i}",
                        filename=filename,
                        file_type="PPTX",
                        slide=slide_num,
                        chunk_index=i,
                        total_chunks=len(text_chunks)
                    )
                    chunks.append(DocumentChunk(content=chunk_text, metadata=metadata))
        
        elif file_ext in ['.xlsx', '.xls']:
            sheets = self.extractor.extract_excel(filepath)
            for text, sheet_name in sheets:
                text_chunks = self.chunker.chunk_text(text)
                for i, chunk_text in enumerate(text_chunks):
                    metadata = DocumentMetadata(
                        chunk_id=f"{filename}_{sheet_name}_c{i}",
                        filename=filename,
                        file_type="EXCEL",
                        sheet=sheet_name,
                        chunk_index=i,
                        total_chunks=len(text_chunks)
                    )
                    chunks.append(DocumentChunk(content=chunk_text, metadata=metadata))
        
        elif file_ext == '.txt':
            text = self.extractor.extract_txt(filepath)
            text_chunks = self.chunker.chunk_text(text)
            for i, chunk_text in enumerate(text_chunks):
                metadata = DocumentMetadata(
                    chunk_id=f"{filename}_c{i}",
                    filename=filename,
                    file_type="TXT",
                    chunk_index=i,
                    total_chunks=len(text_chunks)
                )
                chunks.append(DocumentChunk(content=chunk_text, metadata=metadata))
        
        print(f"  [2/5] Created {len(chunks)} chunks")
        
        # STAGE 2: Generate embeddings
        print(f"  [3/5] Generating embeddings...")
        texts = [chunk.content for chunk in chunks]
        embeddings = self.embedding_model.encode(texts, show_progress_bar=False)
        
        for chunk, embedding in zip(chunks, embeddings):
            chunk.embedding = embedding
        
        print(f"  [4/5] Embeddings generated ({embeddings.shape})")
        
        self.chunks.extend(chunks)
        print(f"  [5/5] ✓ Document processed")
        
        return chunks
    
    def build_indexes(self):
        """Build BM25 and vector indexes"""
        print(f"\nBuilding indexes for {len(self.chunks)} chunks...")
        
        # Build BM25 index
        print("  Building BM25 index...")
        self.bm25_index.index(self.chunks)
        
        # Build FAISS vector index
        print("  Building FAISS vector index...")
        embeddings = np.array([chunk.embedding for chunk in self.chunks])
        
        # Normalize for cosine similarity
        faiss.normalize_L2(embeddings)
        
        # Create FAISS index
        dimension = embeddings.shape[1]
        self.vector_index = faiss.IndexFlatIP(dimension)  # Inner product (cosine after normalization)
        self.vector_index.add(embeddings)
        
        print(f"✓ Indexes built successfully")

